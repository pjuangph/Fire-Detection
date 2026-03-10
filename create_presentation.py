#!/usr/bin/env python3
"""Generate the Fire-Detection presentation (one-time use, then delete)."""

from __future__ import annotations

import os

import yaml
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE


# ── Colour palette ──────────────────────────────────────────────────────
DARK_BG = RGBColor(0x33, 0x33, 0x33)
ACCENT_RED = RGBColor(0xE0, 0x3C, 0x31)
ACCENT_ORANGE = RGBColor(0xF4, 0x8C, 0x06)
ACCENT_BLUE = RGBColor(0x3B, 0x82, 0xF6)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
MED_GRAY = RGBColor(0x99, 0x99, 0x99)
VERY_LIGHT = RGBColor(0xF0, 0xF0, 0xF0)
DARK_TEXT = RGBColor(0x22, 0x22, 0x22)
GREEN = RGBColor(0x2E, 0x7D, 0x32)
GREEN_BG = RGBColor(0xE8, 0xF5, 0xE9)


def _set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_textbox(slide, left, top, width, height, text, font_size=18,
                 color=WHITE, bold=False, alignment=PP_ALIGN.LEFT,
                 font_name='Aptos'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def _add_bullet_frame(slide, left, top, width, height, items,
                      font_size=16, color=WHITE, font_name='Aptos',
                      spacing=Pt(6)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = font_name
        p.level = 0
        p.space_after = spacing
    return txBox


def _add_title_bar(slide, title_text):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
        Inches(13.333), Inches(1.2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_BG
    shape.line.fill.background()
    _add_textbox(slide, Inches(0.6), Inches(0.15), Inches(12), Inches(0.9),
                 title_text, font_size=32, color=WHITE, bold=True)


def _add_subtitle(slide, text, top=Inches(1.35)):
    _add_textbox(slide, Inches(0.6), top, Inches(12), Inches(0.5),
                 text, font_size=18, color=MED_GRAY, bold=False)


def _safe_add_image(slide, path, left, top, width=None, height=None):
    if os.path.isfile(path):
        kwargs = {'left': left, 'top': top}
        if width:
            kwargs['width'] = width
        if height:
            kwargs['height'] = height
        slide.shapes.add_picture(path, **kwargs)
        return True
    return False


def _center_image(slide, path, slide_w=Inches(13.333), slide_h=Inches(7.5)):
    if not os.path.isfile(path):
        return False
    with Image.open(path) as img:
        img_w, img_h = img.size
    aspect = img_w / img_h
    fit_h = slide_h
    fit_w = int(fit_h * aspect)
    if fit_w > slide_w:
        fit_w = slide_w
        fit_h = int(fit_w / aspect)
    left = (slide_w - fit_w) // 2
    top = (slide_h - fit_h) // 2
    slide.shapes.add_picture(path, left, top, fit_w, fit_h)
    return True


def _make_table(slide, rows, left, top, width, height, col_widths=None):
    """Helper: add a styled table and return the table object."""
    tbl = slide.shapes.add_table(
        len(rows), len(rows[0]), left, top, width, height).table
    if col_widths:
        for ci, w in enumerate(col_widths):
            tbl.columns[ci].width = Inches(w)
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = tbl.cell(ri, ci)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(14)
                p.font.name = 'Aptos'
                if ri == 0:
                    p.font.bold = True
                    p.font.color.rgb = WHITE
                else:
                    p.font.color.rgb = DARK_TEXT
                    if ci == 0:
                        p.font.bold = True
            if ri == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = DARK_BG
            elif ri % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = VERY_LIGHT
    return tbl


def _load_yaml_metrics(path):
    """Load metrics from a best_model YAML config, or return None."""
    if not os.path.isfile(path):
        return None
    with open(path) as f:
        cfg = yaml.safe_load(f)
    return cfg.get('metrics')


def _gif_slide(prs, title, gif_path, subtitle=None):
    """Add a full-slide GIF animation with title bar."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, WHITE)
    _add_title_bar(slide, title)
    if subtitle:
        _add_subtitle(slide, subtitle)
    added = _safe_add_image(slide, gif_path,
                            Inches(0.5), Inches(1.5), width=Inches(12.3))
    if not added:
        _add_textbox(slide, Inches(1), Inches(3.5), Inches(11), Inches(1),
                     f'[Animation not found: {gif_path}]',
                     font_size=18, color=MED_GRAY, alignment=PP_ALIGN.CENTER)
    return slide


def _metrics_slide(prs, title, subtitle, yaml_path, convergence_path, model_label):
    """Add a results slide with metrics table + convergence plot."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, WHITE)
    _add_title_bar(slide, title)
    _add_subtitle(slide, subtitle)

    m = _load_yaml_metrics(yaml_path)
    if m:
        rows = [
            ('Metric', 'Value'),
            ('Error Rate', f'{m["error_rate"]:.4f}  (= (FP + FN) / P)'),
            ('True Positives (TP)', f'{m["TP"]:,}'),
            ('False Positives (FP)', f'{m["FP"]:,}'),
            ('False Negatives (FN)', f'{m["FN"]:,}'),
            ('True Negatives (TN)', f'{m["TN"]:,}'),
            ('Precision', f'{m["precision"]:.4f}'),
            ('Recall', f'{m["recall"]:.4f}'),
        ]
        _make_table(slide, rows, Inches(0.8), Inches(2.0),
                    Inches(5.5), Inches(4.0), col_widths=[2.8, 2.7])
    else:
        _add_textbox(slide, Inches(0.8), Inches(3.5), Inches(5.5), Inches(1),
                     f'{model_label} training not yet completed.\n'
                     'Run train-all-models.sh to generate results.',
                     font_size=18, color=MED_GRAY)

    _safe_add_image(slide, convergence_path,
                    Inches(7.0), Inches(1.8), width=Inches(5.8))
    return slide


def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ══════════════════════════════════════════════════════════════════
    # 1 — Title
    # ══════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, WHITE)
    _add_textbox(slide, Inches(1), Inches(1.5), Inches(11), Inches(1.2),
                 'Airborne Fire Detection with Machine Learning',
                 font_size=44, color=DARK_TEXT, bold=True,
                 alignment=PP_ALIGN.CENTER)
    _add_textbox(slide, Inches(1), Inches(3.0), Inches(11), Inches(0.8),
                 'MASTER Hyperspectral Sensor  |  FireSense 2023 Campaign',
                 font_size=24, color=ACCENT_ORANGE, bold=False,
                 alignment=PP_ALIGN.CENTER)
    _add_textbox(slide, Inches(1), Inches(4.2), Inches(11), Inches(0.6),
                 'Multi-Layer Perceptron (MLP)',
                 font_size=20, color=MED_GRAY, alignment=PP_ALIGN.CENTER)
    _add_textbox(slide, Inches(1), Inches(5.2), Inches(11), Inches(0.5),
                 'Team Flaming Kitty: Paht Juangphanich, Codi Lee, Adam Yingling, Candice McDonald',
                 font_size=18, color=MED_GRAY, alignment=PP_ALIGN.CENTER)
    _add_textbox(slide, Inches(1), Inches(6.0), Inches(11), Inches(0.5),
                 'Kaibab Plateau, Arizona  |  October 18\u201320, 2023',
                 font_size=16, color=MED_GRAY, alignment=PP_ALIGN.CENTER)

    # ══════════════════════════════════════════════════════════════════
    # 2 — Problem Statement
    # ══════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, WHITE)
    _add_title_bar(slide, 'Problem Statement')
    _add_bullet_frame(slide, Inches(0.8), Inches(1.8), Inches(5.5), Inches(4.5), [
        'Detect active fire and burn scars from airborne '
        'hyperspectral imagery during prescribed burns',
        'Traditional threshold detectors suffer from false positives '
        '(sun glint, hot soil) and false negatives (smoldering edges)',
        'Goal: reduce error rate (FP + FN) / P while maintaining '
        'near-perfect recall on true fire pixels',
    ], font_size=18, color=DARK_TEXT, spacing=Pt(12))
    _add_textbox(slide, Inches(7.2), Inches(1.8), Inches(5), Inches(0.5),
                 'Key Challenges', font_size=22, color=ACCENT_RED, bold=True)
    _add_bullet_frame(slide, Inches(7.2), Inches(2.5), Inches(5), Inches(4), [
        'Extreme class imbalance: ~1.7% of pixels are fire',
        'Day/night radiometric differences',
        'SWIR saturation near flame front',
        'Pre-burn flight has zero fire \u2014 all detections are FP',
    ], font_size=16, color=DARK_TEXT, spacing=Pt(10))

    # ══════════════════════════════════════════════════════════════════
    # 3 — MASTER Instrument & Dataset
    # ══════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, WHITE)
    _add_title_bar(slide, 'MASTER Instrument & FireSense 2023 Dataset')
    _add_subtitle(slide, '50-channel airborne whisk-broom scanner  |  Prescribed burn, Kaibab Plateau')

    ch_rows = [
        ('Band', 'Channel', '\u03bb', 'Use'),
        ('T4', 'Ch 31', '3.9 \u03bcm', 'Fire detection (MWIR)'),
        ('T11', 'Ch 48', '11.3 \u03bcm', 'Background temp'),
        ('SWIR', 'Ch 22', '2.2 \u03bcm', 'Solar reflection'),
        ('Red', 'Ch 5', '0.65 \u03bcm', 'NDVI'),
        ('NIR', 'Ch 9', '0.87 \u03bcm', 'NDVI'),
    ]
    _make_table(slide, ch_rows, Inches(0.5), Inches(2.0),
                Inches(5.0), Inches(3.0), col_widths=[1.0, 1.0, 1.0, 2.0])

    fl_rows = [
        ('Flight', 'Date', 'Sweeps', 'Condition', 'Fire?'),
        ('03', 'Oct 18', '14', 'Pre-burn (day)', 'No'),
        ('04', 'Oct 19', '22', 'Active fire (night)', 'Yes'),
        ('05', 'Oct 19', '12', 'Overnight (night)', 'Yes'),
        ('06', 'Oct 20', '13', 'Smoldering (day)', 'Yes'),
    ]
    _make_table(slide, fl_rows, Inches(6.0), Inches(2.0),
                Inches(6.8), Inches(3.0), col_widths=[1.0, 1.2, 1.2, 2.2, 1.2])

    _add_bullet_frame(slide, Inches(0.8), Inches(5.5), Inches(11), Inches(1.5), [
        '61 HDF files across 4 flights  |  Each file = one scan sweep (50 channels + lat/lon)',
        'Spatial resolution: ~8 m native pixels \u2192 gridded to 0.00025\u00b0 (\u224828 m)',
    ], font_size=16, color=DARK_TEXT, spacing=Pt(8))

    # ══════════════════════════════════════════════════════════════════
    # 4 — Data Processing: Grid + Multi-Pass + Pipeline
    # ══════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, WHITE)
    _add_title_bar(slide, 'Data Processing Pipeline')
    _add_subtitle(slide, 'From raw HDF scans to ML-ready features via multi-pass gridding')

    steps = [
        ('1. Ingest', 'Read MASTER\nHDF files\n(50 ch)'),
        ('2. Calibrate', 'Planck \u2192\nBrightness\nTemp'),
        ('3. Grid', 'Mosaic onto\nlat/lon grid\n(0.00025\u00b0)'),
        ('4. Accumulate', 'Multi-pass:\nT4_max, dT_max\nobs_count'),
        ('5. Features', '12 aggregate\nfeatures per\ngrid cell'),
        ('6. Detect', 'ML inference\nor threshold\nrules'),
    ]
    box_w, box_h = Inches(1.8), Inches(1.8)
    start_x, y = Inches(0.4), Inches(2.2)
    gap = Inches(0.25)
    for i, (title, body) in enumerate(steps):
        x = start_x + i * (box_w + gap)
        color = ACCENT_RED if i == 5 else DARK_BG
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x, y, box_w, box_h)
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()
        tf = shape.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.font.name = 'Aptos'
        p2 = tf.add_paragraph()
        p2.text = body
        p2.font.size = Pt(11)
        p2.font.color.rgb = LIGHT_GRAY
        p2.font.name = 'Aptos'
        p2.alignment = PP_ALIGN.CENTER
        if i < len(steps) - 1:
            _add_textbox(slide, x + box_w, y + box_h / 2 - Inches(0.15),
                         gap, Inches(0.3), '\u2192',
                         font_size=22, color=ACCENT_ORANGE, bold=True,
                         alignment=PP_ALIGN.CENTER)

    _add_bullet_frame(slide, Inches(0.5), Inches(4.4), Inches(12), Inches(2.8), [
        'Gridding: irregular whisk-broom pixels \u2192 regular lat/lon mosaic  |  '
        '~3 raw pixels per cell',
        'Multi-pass: same cell scanned 1\u201340\u00d7 across sweeps \u2192 '
        'running accumulators track max, sum, count',
        'Multi-pass consistency filter: require fire in \u22652 passes '
        '(52% FP reduction on pre-burn flight)',
        'Hybrid normalization: thermal features (T4, T11, \u0394T) / 573 K '
        '(dry wood ignition); non-thermal features z-scored against pre-burn baseline',
        'Result: one row per grid cell with 12 features \u2192 ready for ML or threshold detection',
    ], font_size=16, color=DARK_TEXT, spacing=Pt(8))

    # ══════════════════════════════════════════════════════════════════
    # 5 — Multi-Pass Scanning (diagram)
    # ══════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, WHITE)
    _center_image(slide, 'plots/diagram_multipass_scanning.png')

    # ══════════════════════════════════════════════════════════════════
    # 6 — Baseline: T4-T11 Threshold + NDVI Vegetation
    # ══════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, WHITE)
    _add_title_bar(slide, 'Baseline Detection: Physics + Vegetation')
    _add_subtitle(slide, 'T4\u2013T11 threshold  |  NDVI vegetation-loss confirmation')

    _add_bullet_frame(slide, Inches(0.8), Inches(2.0), Inches(5.5), Inches(5), [
        'Fire at ~800 K is 23,000\u00d7 brighter than 300 K background at 3.9 \u03bcm',
        'Simple rule: T4 > 325 K AND \u0394T > 10 K',
        'Contextual: T4 > mean + 3\u03c3 of 61\u00d761 neighborhood',
        'SWIR rejection filters solar reflection false alarms',
        '',
        'Vegetation-loss confirmation:',
        'NDVI baseline from first daytime pass (write-once)',
        'Fire confirmed when NDVI drops \u22650.15 at thermal pixel',
        'Confirmed pixels persist as burn scars',
    ], font_size=16, color=DARK_TEXT, spacing=Pt(8))

    _safe_add_image(slide, 'plots/diagram_threshold_explanation.png',
                    Inches(6.5), Inches(1.8), width=Inches(6.5))

    # ══════════════════════════════════════════════════════════════════
    # 7 — ML Approach: MLP Overview
    # ══════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, WHITE)
    _add_title_bar(slide, 'Machine Learning Approach')
    _add_subtitle(slide, 'FireMLP: 12 input features \u2192 P(fire) per grid cell')

    _add_bullet_frame(slide, Inches(0.5), Inches(2.0), Inches(5.8), Inches(5), [
        'Why ML?  Threshold rules miss smoldering edges, '
        'produce FP on hot soil.  ML learns the boundary.',
        '',
        'FireMLP (Multi-Layer Perceptron):',
        '  Trained from scratch on our data',
        '  Variable architecture (grid-searched)',
        '  Compact: ~11K params  |  ~10 KB',
        '  5 loss functions tested across 240 configs',
        '',
        'Two selected models:',
        '  Conservative: minimal false positives on clear days',
        '  Best Overall: lowest total error rate',
    ], font_size=16, color=DARK_TEXT, spacing=Pt(6))

    _safe_add_image(slide, 'plots/diagram_model_comparison.png',
                    Inches(6.5), Inches(1.8), width=Inches(6.5))

    # ══════════════════════════════════════════════════════════════════
    # 8 — What Goes In / What Comes Out (diagram)
    # ══════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, WHITE)
    _center_image(slide, 'plots/diagram_model_io.png')

    # ══════════════════════════════════════════════════════════════════
    # 9 — Loss Functions & Hyperparameter Tuning (all models)
    # ══════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, WHITE)
    _add_title_bar(slide, 'Loss Functions & Hyperparameter Tuning')
    _add_subtitle(slide, '240 MLP configurations tested across 5 loss functions')

    # MLP losses table
    _add_textbox(slide, Inches(0.3), Inches(2.0), Inches(6.0), Inches(0.5),
                 'MLP Grid Search (240 runs)', font_size=20, color=ACCENT_BLUE, bold=True)

    loss_rows = [
        ('Loss Function', 'Runs', 'Best Error Rate'),
        ('Weighted BCE', '108', '0.0511'),
        ('Tversky', '32', '0.0507'),
        ('Error-Rate', '36', '0.0623'),
        ('Focal-Error-Rate', '32', '0.0988'),
        ('Combined', '32', '0.0567'),
    ]
    _make_table(slide, loss_rows, Inches(0.3), Inches(2.6),
                Inches(6.0), Inches(3.0), col_widths=[2.5, 1.0, 2.5])

    # Tuned hyperparameters
    _add_textbox(slide, Inches(6.8), Inches(2.0), Inches(6.0), Inches(0.5),
                 'Hyperparameter Space', font_size=20,
                 color=RGBColor(0x7C, 0x3A, 0xED), bold=True)
    _add_bullet_frame(slide, Inches(6.8), Inches(2.6), Inches(6.0), Inches(3.0), [
        'Architectures: 64\u00d764\u00d764\u00d732, 128\u00d764\u00d732, 128\u00d7128\u00d7128\u00d732/64',
        'Learning rates: 0.01, 0.001, 0.0001',
        'Normalization: standard vs hybrid',
        'Dropout: 0.0 vs 0.1',
        'Importance weights: gt \u2208 {5,10,15}, fire \u2208 {3,5,10}',
        '',
        'Key findings:',
        '  Standard normalization >> hybrid',
        '  No dropout >> dropout 0.1',
        '  Smaller architectures perform better',
    ], font_size=14, color=DARK_TEXT, spacing=Pt(6))

    # Oversampling
    _add_textbox(slide, Inches(0.3), Inches(5.8), Inches(12.5), Inches(0.4),
                 'Handling Class Imbalance', font_size=18, color=ACCENT_RED, bold=True)
    _add_bullet_frame(slide, Inches(0.3), Inches(6.2), Inches(12.5), Inches(1.0), [
        'Only ~1.7% of grid cells contain fire \u2014 minority oversampling balances the training set',
        'Pixel-wise weights: ground-truth flight \u00d710, fire pixels \u00d75, others \u00d71',
    ], font_size=14, color=DARK_TEXT, spacing=Pt(4))

    # ══════════════════════════════════════════════════════════════════
    # 10 — Reading the Real-Time Plots (legend)
    # ══════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, WHITE)
    _center_image(slide, 'plots/diagram_realtime_legend.png')

    # ══════════════════════════════════════════════════════════════════
    # 11 — Two Selected Models: Conservative vs Best Overall
    # ══════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, WHITE)
    _add_title_bar(slide, 'Two Selected Models from 240-Run Grid Search')
    _add_subtitle(slide, 'Conservative (lowest daytime FP) vs Best Overall (lowest error rate)')

    model_rows = [
        ('', 'Conservative (Run #17)', 'Best Overall (Run #156)'),
        ('Loss', 'Weighted BCE', 'Tversky (\u03b1=0.5, \u03b2=0.5)'),
        ('Architecture', '64\u00d764\u00d764\u00d732', '128\u00d764\u00d732'),
        ('Learning Rate', '0.001', '0.001'),
        ('Flight 03 FP', '27', '117'),
        ('Test TP', '8,254', '8,899'),
        ('Test FP', '63', '291'),
        ('Test FN', '814', '169'),
        ('Precision', '99.2%', '96.8%'),
        ('Recall', '91.0%', '98.1%'),
        ('Error Rate', '0.097', '0.051'),
    ]
    _make_table(slide, model_rows, Inches(0.5), Inches(2.0),
                Inches(8.0), Inches(5.0), col_widths=[2.5, 2.75, 2.75])

    _add_bullet_frame(slide, Inches(9.0), Inches(2.5), Inches(4), Inches(4.5), [
        'Conservative: prioritizes zero false alarms on clear-sky daytime imagery',
        '',
        'Best Overall: catches 98% of all fires with lowest total error rate',
        '',
        'Both use standard normalization, no dropout, LR=0.001',
        '',
        'Tradeoff: 27 vs 117 FP on clear day; 91% vs 98% fire recall',
    ], font_size=15, color=DARK_TEXT, spacing=Pt(6))

    # ══════════════════════════════════════════════════════════════════
    # 12-15 — Per-Flight Model Comparison (spatial overlay plots)
    # ══════════════════════════════════════════════════════════════════
    _flight_compare_info = [
        ('2480103', 'Pre-Burn (Flight 03) \u2014 Daytime, No Fire',
         'All detections are false positives  |  Red = Conservative, Purple = Best Overall, Blue = Both'),
        ('2480104', 'Active Fire (Flight 04) \u2014 Night',
         'First fire flight  |  Red = Conservative only, Purple = Best Overall only, Blue = Both'),
        ('2480105', 'Overnight (Flight 05) \u2014 Night',
         'Second night flight  |  Red = Conservative only, Purple = Best Overall only, Blue = Both'),
        ('2480106', 'Daytime Smoldering (Flight 06)',
         'Hardest detection scenario  |  Red = Conservative only, Purple = Best Overall only, Blue = Both'),
    ]
    for fnum_clean, title, subtitle_text in _flight_compare_info:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _set_slide_bg(slide, WHITE)
        _add_title_bar(slide, f'Model Comparison \u2014 {title}')
        _add_subtitle(slide, subtitle_text)
        # Fit image within remaining space (below subtitle, above bottom)
        img_path = f'plots/model_compare_{fnum_clean}.png'
        if os.path.isfile(img_path):
            with Image.open(img_path) as img:
                img_w, img_h = img.size
            aspect = img_w / img_h
            avail_h = Inches(5.5)   # 7.5 - 1.8 top - 0.2 bottom
            avail_w = Inches(12.3)  # 13.333 - margins
            fit_h = avail_h
            fit_w = int(fit_h * aspect)
            if fit_w > avail_w:
                fit_w = avail_w
                fit_h = int(fit_w / aspect)
            left = (Inches(13.333) - fit_w) // 2
            slide.shapes.add_picture(img_path, left, Inches(1.8), fit_w, fit_h)

    # ══════════════════════════════════════════════════════════════════
    # 16 — Model Comparison Summary
    # ══════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, WHITE)
    _add_title_bar(slide, 'Model Comparison Summary')
    _add_subtitle(slide, 'Two selected MLP models from 240-run grid search  |  Evaluated on same test set (4 flights)')

    summary_header = ['', 'Conservative (Run #17)', 'Best Overall (Run #156)']
    summary_rows = [
        summary_header,
        ['Loss Function', 'Weighted BCE', 'Tversky'],
        ['Architecture', '64\u00d764\u00d764\u00d732', '128\u00d764\u00d732'],
        ['Flight 03 FP (daytime)', '27', '117'],
        ['Error Rate', '0.097', '0.051'],
        ['TP', '8,254', '8,899'],
        ['FP', '63', '291'],
        ['FN', '814', '169'],
        ['Precision', '99.2%', '96.8%'],
        ['Recall', '91.0%', '98.1%'],
    ]
    _make_table(slide, summary_rows,
                Inches(0.5), Inches(2.0), Inches(8.5), Inches(4.5),
                col_widths=[3.0, 2.75, 2.75])

    _add_bullet_frame(slide, Inches(9.5), Inches(2.5), Inches(3.5), Inches(4.5), [
        'Conservative: deploy when false alarms are costly (e.g. automated alerts)',
        '',
        'Best Overall: deploy when missing fires is costly (e.g. operational monitoring)',
        '',
        '240 configs tested across 5 loss functions',
        '',
        'Error Rate = (FP + FN) / P',
    ], font_size=14, color=DARK_TEXT, spacing=Pt(6))

    # ══════════════════════════════════════════════════════════════════
    # 20 — Conclusion + Future Work
    # ══════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, WHITE)
    _add_title_bar(slide, 'Conclusion & Future Work')

    # Best models callout
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(2.0),
        Inches(6), Inches(3.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = GREEN_BG
    shape.line.color.rgb = GREEN
    shape.line.width = Pt(2)

    _add_textbox(slide, Inches(0.8), Inches(2.1), Inches(5.4), Inches(0.5),
                 '\u2b50 Two Production Models from 240-Run Search',
                 font_size=20, color=GREEN, bold=True)
    _add_bullet_frame(slide, Inches(0.8), Inches(2.7), Inches(5.4), Inches(2.5), [
        'Conservative (Run #17, BCE):',
        '  27 FP on clear day  |  Precision 99.2%  |  Recall 91.0%',
        '',
        'Best Overall (Run #156, Tversky):',
        '  Error rate 0.051  |  Precision 96.8%  |  Recall 98.1%',
        '',
        'Compact networks (<10 KB) train in seconds',
    ], font_size=14, color=DARK_TEXT, spacing=Pt(4))

    # Future work
    _add_textbox(slide, Inches(7.0), Inches(2.0), Inches(5.5), Inches(0.5),
                 'Next Steps', font_size=22, color=ACCENT_BLUE, bold=True)
    _add_bullet_frame(slide, Inches(7.0), Inches(2.7), Inches(5.5), Inches(4.5), [
        'Ensemble: combine conservative + best overall predictions',
        'Transfer learning to new fire campaigns',
        'Spatial context features (neighbor stats)',
        'Temporal features from sequential passes',
        'Operational deployment in real-time pipeline',
        'Cross-validation across flights',
    ], font_size=16, color=DARK_TEXT, spacing=Pt(8))

    # ══════════════════════════════════════════════════════════════════
    # 21 — Summary
    # ══════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, WHITE)
    _add_textbox(slide, Inches(1), Inches(1.5), Inches(11), Inches(1),
                 'Summary', font_size=40, color=DARK_TEXT, bold=True,
                 alignment=PP_ALIGN.CENTER)
    _add_bullet_frame(slide, Inches(2), Inches(3.0), Inches(9), Inches(3.5), [
        'Physics-based fire detection from 50-channel MASTER hyperspectral data',
        'Multi-pass consistency + vegetation-loss confirmation reduce false positives',
        '240 MLP configurations tested across 5 loss functions',
        'Conservative model: 27 FP on clear day (99.2% precision)',
        'Best Overall model: 0.051 error rate with 98.1% recall',
    ], font_size=20, color=DARK_TEXT, spacing=Pt(14))
    _add_textbox(slide, Inches(1), Inches(6.5), Inches(11), Inches(0.5),
                 'github.com  |  FireSense 2023  |  MASTER L1B',
                 font_size=14, color=MED_GRAY, alignment=PP_ALIGN.CENTER)

    # ══════════════════════════════════════════════════════════════════
    # BACKUP SLIDES
    # ══════════════════════════════════════════════════════════════════

    # Backup separator
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, VERY_LIGHT)
    _add_textbox(slide, Inches(1), Inches(3), Inches(11), Inches(1),
                 'Backup Slides', font_size=40, color=MED_GRAY, bold=True,
                 alignment=PP_ALIGN.CENTER)

    # B1 — MLP vs Baseline: Active Fire (Night)
    _gif_slide(prs,
               'Backup: MLP vs Baseline \u2014 Active Fire Night (Flight 04)',
               'plots/gifs/compare_ml_vs_simple_2480104_active.gif')

    # B2 — MLP vs Baseline: Overnight
    _gif_slide(prs,
               'Backup: MLP vs Baseline \u2014 Overnight (Flight 05)',
               'plots/gifs/compare_ml_vs_simple_2480105_overnight.gif')

    # (TabPFN backup slides removed)

    # B5 — MLP Prediction Maps
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, WHITE)
    _add_title_bar(slide, 'Backup: MLP Spatial Prediction Maps')
    _safe_add_image(slide, 'plots/tune_prediction_map_2480106.png',
                    Inches(1.5), Inches(1.8), width=Inches(10))

    # B6 — MLP Probability Calibration
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, WHITE)
    _add_title_bar(slide, 'Backup: MLP Probability Calibration')
    _safe_add_image(slide, 'plots/tune_probability_hist.png',
                    Inches(3), Inches(1.8), width=Inches(7))

    # B7 — Grid Resolution Comparison
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, WHITE)
    _add_title_bar(slide, 'Backup: Grid Resolution Comparison')
    _safe_add_image(slide, 'plots/grid_resolution_comparison.png',
                    Inches(0.5), Inches(1.8), width=Inches(12))

    # ── Disable auto-advance on all slides (manual click only) ──
    from lxml import etree
    nsmap = {'p': 'http://schemas.openxmlformats.org/presentationml/2006/main'}
    for slide in prs.slides:
        # Remove any existing transition element
        for trans in slide._element.findall('p:transition', nsmap):
            slide._element.remove(trans)
        # Add transition with advClick=true, advTm removed (no auto-advance)
        trans_el = etree.SubElement(
            slide._element,
            '{http://schemas.openxmlformats.org/presentationml/2006/main}transition')
        trans_el.set('advClick', '1')

    # ── Save ──
    out_path = 'Fire_Detection_Presentation.pptx'
    prs.save(out_path)
    n_main = 18
    n_total = len(prs.slides)
    print(f'Saved presentation to {out_path}')
    print(f'  {n_main} main slides + {n_total - n_main} backup = {n_total} total')


if __name__ == '__main__':
    create_presentation()
